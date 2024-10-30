# pip install gradio==4.44.1

import os
os.environ["OPENAI_API_KEY"] = "sk-XXXXXXXX"
os.environ['ANTHROPIC_API_KEY'] = "sk-XXXXXX"
import logging
import sys
import time
import base64
import concurrent
import re
import json
import tempfile
import gradio as gr
import warnings
import random
import threading
import openpyxl
import glob
import base64
import anthropic

from llama_index.graph_stores.neo4j import Neo4jPropertyGraphStore
from llama_index.core import PropertyGraphIndex
from llama_index.embeddings.openai import OpenAIEmbedding
from llama_index.llms.openai import OpenAI as LlamaOpenAI

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from pdf2image import convert_from_path
from pdf2image.exceptions import (
    PDFInfoNotInstalledError,
    PDFPageCountError,
    PDFSyntaxError
)
from pdfminer.high_level import extract_text
from io import BytesIO
from tqdm import tqdm
from openai import OpenAI
from llama_index.core import Document

from docx.oxml.ns import qn
from docx.parts.image import ImagePart
from lxml import etree
import sympy as sp

from concurrent.futures import ThreadPoolExecutor, as_completed

# logging.basicConfig(
#     stream=sys.stdout, level=logging.ERROR
# )  # logging.DEBUG for more verbose output

warnings.filterwarnings("ignore")
# 设置 neo4j 日志记录级别为 ERROR 这样可以减少不必要的日志输出
logging.getLogger("neo4j").setLevel(logging.ERROR)

client = OpenAI()
client_claude = anthropic.Anthropic(
    # This is the default and can be omitted
    api_key=os.getenv("ANTHROPIC_API_KEY"),
)
# 设定最大线程数max_workers_number 如果提高线程数，需要调整neo4j的配置文件 提高资源使用
max_workers_number = 20
# max_workers_number_for_claude 指定claude的线程数，这个需要key的tier支持，否则会报429错误
max_workers_number_for_claude = 2
batch_size = 5
token_lock = threading.Lock()
token_counts = {
            'input': 0,
            'output': 0,
            'cache_read': 0,
            'cache_creation': 0
        }

failed_sections = []
upload_files_path=[]

graph_store = Neo4jPropertyGraphStore(
    username="neo4j",
    password="",
    url="neo4j://localhost:7687",
    database="neo4j",
    )

index_query = PropertyGraphIndex.from_existing(
    property_graph_store=graph_store,
    llm=LlamaOpenAI(model="gpt-4o-mini", temperature=0),
    embed_model=OpenAIEmbedding(model_name="text-embedding-ada-002"),
)

def convert_doc_to_images(pdf_path):
    images = convert_from_path(pdf_path)
    return images

def extract_text_from_doc(pdf_path):
    text = extract_text(pdf_path)
    page_text = []
    return text

def pdf_analysis(pdf_path):
    images = convert_doc_to_images(pdf_path)
    text = extract_text_from_doc(pdf_path)
    # for img in images:
    #     display(img)
    return images, text

# Converting images to base64 encoded images in a data URI format to use with the ChatCompletions API
def get_img_uri(img):
    buffer = BytesIO()
    img.save(buffer, format="jpeg")
    base64_image = base64.b64encode(buffer.getvalue()).decode("utf-8")
    data_uri = f"data:image/jpeg;base64,{base64_image}"
    return data_uri

def analyze_image(img_url):
    system_prompt = '''
        你是一个非常有帮助的AI助手，你的任务是帮助用户分析图片中的信息。
        
        根据提供的PDF的image，需要你分析这张图片并提取出图片中的所有信息，然后输出。注意不可以省略任何信息。省略信息将会导致用户无法理解图片中的内容。
        
        所有最终输出结果都放在<results> tags标签中。
        
        不需要输出中间分析过程，只需要输出最终结果。
    '''

    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {
                "role": "system",
                "content": system_prompt
            },
            {
                "role": "user",
                "content": [
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": f"{img_url}"
                        }
                    },
                ]
            },
        ],
        temperature=0,
        max_tokens=6000,
        top_p=0.95,
        response_format={
        "type": "text"
        },
        logprobs=True
    )

    # 获取输入和输出 token 数量
    print(f"提示使用的token数量: {response.usage.prompt_tokens}")
    print(f"回复生成的token数量: {response.usage.completion_tokens}")
    print(f"总共使用的token数量: {response.usage.total_tokens}")
    input_tokens = response.usage.prompt_tokens
    output_tokens = response.usage.completion_tokens
    return response.choices[0].message.content,input_tokens,output_tokens

def analyze_doc_image(img):
    img_uri = get_img_uri(img)
    interpreter_result,input_tokens,output_tokens = analyze_image(img_uri)
    # 提取tag对象内容
    try:
        pattern = r"<results>(.*?)</results>"
        match = re.search(pattern, interpreter_result, re.DOTALL)  # 使用re.DOTALL标志
        if match:
            interpreter_result = match.group(1)
    except Exception as e:
        print("Error in extracting tag content:", e)
    return interpreter_result

# 提取json内容返回整理后的数据list
def pdf_adjust_split_json_contents(json_contents):
    documents=[]
    sections = []
    # 提取并打印每个文件名和对应的页面描述
    for item in json_contents:
        file_path = item.get("file_path")
        pages_description = item.get("pages_description", [])
        print(f"文件名: {file_path}")
        for content in pages_description:
            file_sections={}
            file_sections["file_path"] = file_path
            # 使用双换行符（空行）作为分割标志
            # file_sections["content"] = content.split('\n\n')
            file_sections["content"] = content
            sections.append(file_sections)
        print(f"文件中的页面数量: {len(pages_description)}处理完成。")
    # 使用列表推导式去掉空字符串
    cleaned_list = [item for item in sections if item["content"] != ""]
    for idx,section in enumerate(cleaned_list):
        # 将 section 中的所有反斜杠 \ 替换为正斜杠 /
        sub_section = section["content"].strip().replace('\\', '/')
        documents.append(Document(text=sub_section, id_=f"doc_id_{idx}", metadata={"file_path": section["file_path"]}))
        print(f"section {idx + 1}：")
        print(section["file_path"])
        print("_id:",f"doc_id_{idx}")
        print(sub_section)
    print("==========len documents:",len(documents))
    return documents

# 提取json内容返回整理后的数据list
def docx_adjust_split_json_contents(json_contents):
    documents=[]
    # 提取并打印每个文件名和对应的页面描述
    for idx, item in enumerate(json_contents):
        file_path = item.get("docx_path")
        pages_description = item.get("content", [])
        if pages_description is None or pages_description == "" or pages_description.strip() == "":
            continue
        # 将 pages_description 中的所有反斜杠 \ 替换为正斜杠 /
        pages_description = pages_description.strip().replace('\\', '/')
        documents.append(Document(text=pages_description, id_=f"doc_id_{idx}", metadata={"file_path": file_path, "images": item.get("images", [])}))
        print(f"pages_description {idx + 1}：")
        print("_id:",f"doc_id_{idx}")
        print(pages_description)
    print("==========len documents:",len(documents))
    return documents

# 提取json内容返回整理后的数据list
def txt_adjust_split_json_contents(json_contents):
    documents=[]
    # 提取并打印每个文件名和对应的页面描述
    for idx, item in enumerate(json_contents):
        file_path = item.get("file_path")
        pages_description = item.get("original_content", [])+"\n\n"+item.get("contextualized_content", [])
        if pages_description is None or pages_description == "" or pages_description.strip() == "":
            continue
        # 将 pages_description 中的所有反斜杠 \ 替换为正斜杠 /
        pages_description = pages_description.strip().replace('\\', '/')
        documents.append(Document(text=pages_description, id_=f"doc_id_{idx}", metadata={"file_path": file_path}))
        print(f"pages_description {idx + 1}：")
        print("_id:",f"doc_id_{idx}")
        print(pages_description)
    print("==========len documents:",len(documents))
    return documents

# 提取json内容返回整理后的数据list
def pptx_adjust_split_json_contents(json_contents):
    documents=[]
    # 提取并打印每个文件名和对应的页面描述
    for idx, item in enumerate(json_contents):
        file_path = item.get("file_path")
        pages_description = item.get("content", [])
        if pages_description is None or pages_description == "" or pages_description.strip() == "":
            continue
        # 将 pages_description 中的所有反斜杠 \ 替换为正斜杠 /
        pages_description = pages_description.strip().replace('\\', '/')
        documents.append(Document(text=pages_description, id_=f"doc_id_{idx}", metadata={"file_path": file_path, "images": item.get("images", [])}))
        print(f"pages_description {idx + 1}：")
        print("_id:",f"doc_id_{idx}")
        print(pages_description)
    print("==========len documents:",len(documents))
    return documents

def insert_section(batch):
    global failed_sections

    index = PropertyGraphIndex.from_existing(
        property_graph_store=graph_store,
        llm=LlamaOpenAI(model="gpt-4o-mini", temperature=0),
        embed_model=OpenAIEmbedding(model_name="text-embedding-ada-002"),
    )

    for document in batch:
        try:
            # 把内容插入到 Graph 数据库中
            index.insert(document)
            print(f"Inserted document: {document.text[:50]} in thread {threading.current_thread().name}...")
        except TimeoutError as te:
            print(f"Thread {threading.current_thread().name} timeout error: {te}")
            if document not in failed_sections:
                failed_sections.append(document)
        except Exception as e:
            print(f"Error inserting document: {e}  section: {document.text[:80]}...")
            if document not in failed_sections:
                failed_sections.append(document)
        finally:
            print(f"Thread {threading.current_thread().name} has finished.")

# 并行处理文档插入
def parallel_insert(documents_result, batch_size=batch_size, max_workers_number=max_workers_number):
    total_batches = len(documents_result) // batch_size + (1 if len(documents_result) % batch_size != 0 else 0)
    print(f"Total batches: {total_batches}, batch size: {batch_size}, max workers: {max_workers_number}")

    futures = []
    documents=[]
    for i in range(total_batches):
        start = i * batch_size
        end = start + batch_size
        documents.append(documents_result[start:end])

    with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers_number) as executor:
        futures = [executor.submit(insert_section, batch) for batch in documents]

    for future in concurrent.futures.as_completed(futures):
        try:
            future.result()  # 检查线程执行结果，捕获异常
        except concurrent.futures.TimeoutError:
            print(f"Thread task exceeded the timeout limit and was cancelled.")
        except Exception as e:
            print(f"Thread generated an exception: {e}")

def pdf_write_to_graphRAG(history,pdf_json_path):
    # 如果文件不存在或发生其他错误，将打印错误信息并终止程序。
    try:
        pdf_json_path = os.path.join(os.getcwd(), pdf_json_path, "pdf_analyzed_json.json")
        # 打开并加载 JSON 文件
        with open(pdf_json_path, 'r', encoding='utf-8') as f:
            json_contents = json.load(f)
    except FileNotFoundError:
        print(f"错误: 文件 '{pdf_json_path}' 不存在。")
        history.append(["", f"错误: 文件 '{pdf_json_path}' 不存在。"])
        return history
    except json.JSONDecodeError as e:
        print(f"错误: JSON解析失败 - {e}")
        history.append(["", f"错误: JSON解析失败 - {e}"])
        return history
    except Exception as e:
        print(f"发生未知错误: {e}")
        history.append(["", f"发生未知错误: {e}"])
        return history

    # 测量函数执行时间
    start_time = time.time()

    index = PropertyGraphIndex.from_existing(
        property_graph_store=graph_store,
        llm=LlamaOpenAI(model="gpt-4o-mini", temperature=0),
        embed_model=OpenAIEmbedding(model_name="text-embedding-ada-002"),
    )

    # 提取json内容返回整理后的数据list
    documents_result = pdf_adjust_split_json_contents(json_contents)
    # 调用并行插入函数，默认使用N个并行线程
    parallel_insert(documents_result)
    print("====================Failed sections processing again====================")
    print("Failed to insert the following sections:",failed_sections)
    if len(failed_sections)>0 and failed_sections:
        print("Failed to insert the following sections:")
        for section in failed_sections:
            try:
                # 把内容插入到 Graph 数据库中
                document = Document(text=section.text ,metadata={"file_path": section.metadata['file_path']})
                index.insert(document)
                print(f"Inserted document: {section.text[:50]}...")
            except Exception as e:
                print(f"Error inserting document: {e}  section: {section.text[:80]}...")
    # 计算并打印执行时间
    execution_time = time.time() - start_time
    print(f"数据导入neo4j执行时间: {round(execution_time, 2)} 秒")
    print("数据导入neo4j完成")
    history.append(["", f"全部PDF数据导入完成。全部数据导入执行时间: {round(execution_time, 2)} 秒"])
    return history

def docx_write_to_graphRAG(history,docx_json_path):
    if not os.path.isdir(os.path.join(os.getcwd(), docx_json_path)):
        print(f"错误: '{os.path.join(os.getcwd(), docx_json_path)}' 不是一个有效的目录。")
        history.append(["", f"错误: '{os.path.join(os.getcwd(), docx_json_path)}' 不是一个有效的目录。"])
        return history

    all_json_contents = {}
    for file_name in os.listdir(os.path.join(os.getcwd(), docx_json_path)):
        if file_name.lower().endswith('.json'):
            file_path = os.path.join(os.getcwd(), docx_json_path, file_name)
            # 如果文件不存在或发生其他错误，将打印错误信息并终止程序。
            try:
                # 打开并加载 JSON 文件
                with open(file_path, 'r', encoding='utf-8') as f:
                    json_contents = json.load(f)
                    all_json_contents[file_name] = json_contents
            except FileNotFoundError:
                print(f"错误: 文件 '{file_path}' 不存在。")
                history.append(["", f"错误: 文件 '{file_path}' 不存在。"])
                return history
            except json.JSONDecodeError as e:
                print(f"错误: JSON解析失败 - {e}")
                history.append(["", f"错误: JSON解析失败 - {e}"])
                return history
            except Exception as e:
                print(f"发生未知错误: {e}")
                history.append(["", f"发生未知错误: {e}"])
                return history

    index = PropertyGraphIndex.from_existing(
        property_graph_store=graph_store,
        llm=LlamaOpenAI(model="gpt-4o-mini", temperature=0),
        embed_model=OpenAIEmbedding(model_name="text-embedding-ada-002"),
    )

    total_start_time = time.time()
    # 提取json内容返回整理后的数据list
    for file_name, json_contents in all_json_contents.items():
        # 测量函数执行时间
        start_time = time.time()
        documents_result = docx_adjust_split_json_contents(json_contents)
        # 调用并行插入函数，默认使用N个并行线程
        parallel_insert(documents_result)
        print("====================Failed sections processing again====================")
        if failed_sections:
            print("Failed to insert the following sections:")
            for document in failed_sections:
                try:
                    # 把内容插入到 Graph 数据库中
                    index.insert(document)
                    print(f"Inserted document: {document.text[:50]}...")
                except Exception as e:
                    print(f"Error inserting document: {e}  section: {document.text[:80]}...")
        # 计算并打印执行时间
        execution_time = time.time() - start_time
        print(f"{file_name}数据导入neo4j执行时间: {round(execution_time, 2)} 秒")
        print(f"{file_name}数据导入neo4j完成")
    # 计算并打印执行时间
    execution_time = time.time() - total_start_time
    print(f"全部数据导入neo4j执行时间: {round(execution_time, 2)} 秒")
    print("全部数据导入neo4j完成")
    history.append(["", f"全部Word数据导入完成。全部数据导入执行时间: {round(execution_time, 2)} 秒"])
    return history

# 上下文内容提取操作
def process_chunk(doc, chunk, file_path):
    global token_counts
    #for each chunk, produce the context
    contextualized_text, usage = situate_context(doc, chunk)
    with token_lock:
        token_counts['input'] += usage.input_tokens
        token_counts['output'] += usage.output_tokens
        token_counts['cache_read'] += usage.cache_read_input_tokens
        token_counts['cache_creation'] += usage.cache_creation_input_tokens

    print(f"============= token_input: {token_counts['input']}\ntoken_output: {token_counts['output']}\ntoken_cache_read: {token_counts['cache_read']}\ntoken_cache_creation: {token_counts['cache_creation']}")
    return {
        #append the context to the original text chunk
        'text_to_embed': f"{chunk}\n\n{contextualized_text}",
        'metadata': {
            'file_path': file_path,
            'original_content': chunk,
            'contextualized_content': contextualized_text
        }
    }

def txt_write_to_graphRAG(history,txt_json_path):
    if not os.path.isdir(txt_json_path):
        print(f"错误: '{txt_json_path}' 不是一个有效的目录。")
        history.append(["", f"错误: '{txt_json_path}' 不是一个有效的目录。"])
        return history

    all_json_contents = {}
    for file_name in os.listdir(txt_json_path):
        if file_name.lower().endswith('.json'):
            file_path = os.path.join(txt_json_path, file_name)
            # 如果文件不存在或发生其他错误，将打印错误信息并终止程序。
            try:
                # 打开并加载 JSON 文件
                with open(file_path, 'r', encoding='utf-8') as f:
                    json_contents = json.load(f)
                    all_json_contents[file_name] = json_contents
            except FileNotFoundError:
                print(f"错误: 文件 '{file_path}' 不存在。")
                history.append(["", f"错误: 文件 '{file_path}' 不存在。"])
            except json.JSONDecodeError as e:
                print(f"错误: JSON解析失败 - {e}")
                history.append(["", f"错误: JSON解析失败 - {e}"])
            except Exception as e:
                print(f"发生未知错误: {e}")
                history.append(["", f"发生未知错误: {e}"])
    # 上下文内容提取操作
    for file_name, json_contents in all_json_contents.items():
        print(f"===================文件名: {file_name}")
        texts_to_embed = []
        metadata = []
        futures = []
        total_chunks = len(json_contents)
        with open(json_contents[0]['file_path'], 'r', encoding='utf-8') as f:
            temp_all_content = f.read()
        # print("=============================\n{}\n=============================\nfinish\n".format(temp_all_content))
        print(f"Processing {total_chunks} chunks with {max_workers_number_for_claude} threads")
        with ThreadPoolExecutor(max_workers=max_workers_number_for_claude) as executor:
            futures = [executor.submit(process_chunk, temp_all_content, doc['content'], doc['file_path']) for doc in json_contents]

            for future in tqdm(as_completed(futures), total=total_chunks, desc="Processing chunks"):
                try:
                    result = future.result()
                    # print(f"=====================2",result['text_to_embed'])
                    texts_to_embed.append(result['text_to_embed'])
                    metadata.append(result['metadata'])
                except Exception as e:
                    print(f"Thread generated an exception: {e}")

        # 添加了上下文内容的dict类型list数据metadata重新写入all_json_contents变量中
        all_json_contents[file_name] = metadata
        # 把metadata内容写到json文件中
        with open(file_path, 'w', encoding='utf-8') as f:
            json.dump(metadata, f, ensure_ascii=False, indent=4)
        print(f"文件 '{file_name}' 写入完成。")

    index = PropertyGraphIndex.from_existing(
        property_graph_store=graph_store,
        llm=LlamaOpenAI(model="gpt-4o-mini", temperature=0),
        embed_model=OpenAIEmbedding(model_name="text-embedding-ada-002"),
    )

    total_start_time = time.time()
    # 提取json内容返回整理后的数据list
    for file_name, json_contents in all_json_contents.items():
        # 测量函数执行时间
        start_time = time.time()
        documents_result = txt_adjust_split_json_contents(json_contents)
        # 调用并行插入函数，默认使用N个并行线程
        parallel_insert(documents_result)
        print("====================Failed sections processing again====================")
        if failed_sections:
            print("Failed to insert the following sections:")
            for document in failed_sections:
                try:
                    # 把内容插入到 Graph 数据库中
                    index.insert(document)
                    print(f"Inserted document: {document.text[:50]}...")
                except Exception as e:
                    print(f"Error inserting document: {e}  section: {document.text[:80]}...")
        # 计算并打印执行时间
        execution_time = time.time() - start_time
        print(f"{file_name}数据导入neo4j执行时间: {round(execution_time, 2)} 秒")
        print(f"{file_name}数据导入neo4j完成")
    # 计算并打印执行时间
    execution_time = time.time() - total_start_time
    print(f"全部数据导入neo4j执行时间: {round(execution_time, 2)} 秒")
    print("全部数据导入neo4j完成")
    history.append(["", f"全部Txt, Py数据导入完成。全部数据导入执行时间: {round(execution_time, 2)} 秒"])
    return history

# 提取pdf内容并保存到json文件
def extract_pdf_contents(files_path,pdf_json_path):
    docs = []
    for pdf_file in files_path:
        print(pdf_file)
        pages_description=[]
        doc_dict = {}
        images, text=pdf_analysis(pdf_file)
        # Concurrent execution
        with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers_number) as executor:
            futures = [
                executor.submit(analyze_doc_image, img)
                for img in images[0:]
            ]
            
            with tqdm(total=len(images)-1) as pbar:
                for _ in concurrent.futures.as_completed(futures):
                    pbar.update(1)
            
            for f in futures:
                res = f.result()
                pages_description.append(res)
        doc_dict['file_path'] = pdf_file
        doc_dict['pages_description'] = pages_description
        docs.append(doc_dict)
    
    pdf_json_path=os.path.join(os.getcwd(),pdf_json_path,"pdf_analyzed_json.json")
    if not os.path.exists(pdf_json_path):
        os.makedirs(os.path.dirname(pdf_json_path), exist_ok=True)
    # Saving result to file for later
    with open(pdf_json_path, 'w', encoding='utf-8') as f:
        json.dump(docs, f, ensure_ascii=False, indent=4)
    print(f'PDF文件内容提取完成，已保存到文件: {pdf_json_path}')

def parse_formula(element):
    # 递归解析公式中的所有元素
    if element.tag.endswith('oMath'):
        print('oMath:', element.text)
    for child in element:
        parse_formula(child)

def extract_formulas(doc_path):
    document = Document(doc_path)
    formulas = []
    # 遍历文档中的段落
    for paragraph in document.paragraphs:
        for run in paragraph.runs:
            if run.element.xpath('.//m:oMath'):
                # 发现公式
                for element in run.element.xpath('.//m:oMath'):
                    formulas.append(element)
                    parse_formula(element)
    return formulas

def convert_to_latex(formula_elements):
    latex_formulas = []
    for element in formula_elements:
        # 将公式元素转换为LaTeX格式（示例代码，此处需要实际实现）
        # 示例：假设所有公式为 'x^2 + y^2'
        expr = sp.sympify('x^2 + y^2')
        latex_code = sp.latex(expr)
        latex_formulas.append(latex_code)
    return latex_formulas

def parse_formula(element):
    # 递归解析公式中的所有元素
    if element.tag.endswith('oMath'):
        print('oMath:', element.text)
    for child in element:
        parse_formula(child)

def convert_to_latex(formula_elements):
    latex_formulas = []
    for element in formula_elements:
        # 将公式元素转换为LaTeX格式（示例代码，此处需要实际实现）
        # 示例：假设所有公式为 'x^2 + y^2'
        expr = sp.sympify('x^2 + y^2')
        latex_code = sp.latex(expr)
        latex_formulas.append(latex_code)
    return latex_formulas

# 获取图片（该行只能有一个图片）
def get_ImagePart(graph,doc):
    images = graph._element.xpath('.//pic:pic')  # 获取所有图片
    for image in images:
        for img_id in image.xpath('.//a:blip/@r:embed'):  # 获取图片id
            part = doc.part.related_parts[img_id]  # 根据图片id获取对应的图片
            if isinstance(part, ImagePart):
                return part
    return None

# 该行只能有一个图片
def is_image(graph,doc):
    images = graph._element.xpath('.//pic:pic')  # 获取所有图片
    for image in images:
        for img_id in image.xpath('.//a:blip/@r:embed'):  # 获取图片id
            part = doc.part.related_parts[img_id]  # 根据图片id获取对应的图片
            if isinstance(part, ImagePart):
                return True
    return False

# word文档内容提取
def extract_docx_contents(docx_paths, docx_json_path, split_lines=10):
    from docx import Document
    # 检查父目录是否存在，如果不存在就创建
    if not os.path.exists(docx_json_path):
        os.makedirs(docx_json_path)
        print(f"创建输出文件夹: {docx_json_path}")
    for docx_path in docx_paths:
        document = Document(docx_path)
        file_name = os.path.basename(docx_path)
        number = 0
        temp_list = []
        docs = []
        new_docs = []
        # 遍历文档中的段落 这个方法最终只能提取文本内容，图片等无法提取
        for paragraph in document.paragraphs:
            temp = {}
            temp["docx_path"] = docx_path
            temp["content"] = paragraph.text
            imgs = []
            if is_image(paragraph, document):
                # 生成图片path链接
                docx_file_name = os.path.basename(docx_path)
                img_save_path = f"{docx_file_name}_img_{number}.png"
                img_save_path = os.path.join(os.path.dirname(docx_path), img_save_path)
                # print("图片内容保存到：", img_save_path)
                imgs.append(img_save_path)
                temp_return = get_ImagePart(paragraph, document)
                if temp_return:
                    with open(img_save_path, 'wb') as f:
                        f.write(temp_return.blob)  # 将图片的二进制数据写入文件
                    # print('图片保存成功')
                    number += 1
            # 遍历文档中的所有公式
            formula_elements = []
            for run in paragraph.runs:
                if run.element.xpath('.//m:oMath'):
                    # 发现公式
                    for element in run.element.xpath('.//m:oMath'):
                        formula_elements.append(element)
                        parse_formula(element)
            latex_formulas = convert_to_latex(formula_elements)
            # 如果latex_formulas不是空的话，就添加到temp["content"]中
            if latex_formulas:
                temp["content"] += str(latex_formulas)
                print('LaTeX Formulas:', latex_formulas)
            if (temp["content"] is None or temp["content"] == "" or temp["content"].strip() == "") and len(imgs) == 0:
                continue
            temp["images"] = imgs
            temp_list.append(temp)

        # 计算每一个item里面字符数的平均值
        total_chars = 0
        for item in temp_list:
            total_chars += len(item["content"])
        avg_chars = total_chars // len(temp_list)
        print(f"每个item的平均字符数: {avg_chars}")
        
        # 合并N行内容到一个content中
        merged_docs = []
        for i in range(0, len(temp_list), split_lines):
            merged_content = ""
            merged_images = []
            current_docx_path = temp_list[i]["docx_path"]
            for j in range(i, min(i + split_lines, len(temp_list))):
                if temp_list[j]["docx_path"] != current_docx_path:
                    # 如果docx_path不一致，先保存当前合并的内容，然后重新开始合并
                    if merged_content:
                        merged_docs.append({
                            "docx_path": current_docx_path,
                            "content": merged_content.strip(),
                            "images": merged_images
                        })
                    merged_content = ""
                    merged_images = []
                    current_docx_path = temp_list[j]["docx_path"]
                merged_content += temp_list[j]["content"] + "\n"
                merged_images.extend(temp_list[j]["images"])
            if merged_content:
                merged_docs.append({
                    "docx_path": current_docx_path,
                    "content": merged_content.strip(),
                    "images": merged_images
                })
        docs.extend(merged_docs)
    
        for temp in docs:
            if (temp["content"] is None or temp["content"] == "" or temp["content"].strip() == "") and len(temp["images"])==0:
                    continue
            if (temp["content"] is None or temp["content"] == "" or temp["content"].strip() == "") and len(temp["images"])>0:
                temp["content"]="这里是图片信息："+str(temp["images"][0])
            new_docs.append(temp)

        # print("\n\n")
        # print(new_docs)
        # Saving result to file for later
        path = os.path.join(docx_json_path, file_name+".json")
        with open(path, 'w', encoding='utf-8') as f:
            json.dump(new_docs, f, ensure_ascii=False, indent=4)
        print(f"{file_name} Word文件内容提取完成，已保存到文件: {path}")

def situate_context(doc: str, chunk: str):
    DOCUMENT_CONTEXT_PROMPT = """
    <document>
    {doc_content}
    </document>
    """

    CHUNK_CONTEXT_PROMPT = """
    Here is the chunk we want to situate within the whole document
    <chunk>
    {chunk_content}
    </chunk>

    Please give a short succinct context to situate this chunk within the overall document for the purposes of improving search retrieval of the chunk.
    Answer only with the succinct context and nothing else.
    注意：必须使用中文回答并输出结果。
    """
    while True:
        try:
            response = client_claude.messages.create(
                model="claude-3-haiku-20240307",
                max_tokens=2000,
                temperature=0.0,
                messages=[
                    {
                        "role": "user", 
                        "content": [
                            {
                                "type": "text",
                                "text": DOCUMENT_CONTEXT_PROMPT.format(doc_content=doc),
                                "cache_control": {"type": "ephemeral"}  # we will make use of prompt caching for the full documents
                            },
                            {
                                "type": "text",
                                "text": CHUNK_CONTEXT_PROMPT.format(chunk_content=chunk),
                            },
                        ]
                    },
                ],
                extra_headers={"anthropic-beta": "prompt-caching-2024-07-31"}
            )
            print("=======================1", response.content[0].text)
            return response.content[0].text, response.usage  # 成功时直接返回并退出循环
        except Exception as e:
            print(f"Thread generated an exception: {e}")
            if '429' in str(e):  # 检查是否为429错误
                print("遇到429错误，等待80秒后重试...")
                time.sleep(80)  # 等待80秒
            else:
                break  # 非429错误退出循环

# TXT文档内容提取
def extract_txt_contents(txt_files,txt_json_path,split_lines=10):
    # 检查父目录是否存在，如果不存在就创建
    if not os.path.exists(txt_json_path):
        os.makedirs(txt_json_path)
        print(f"创建输出文件夹: {txt_json_path}")
        
    # ======================================================
    # 由于模型限制，每个item的字符数不能超过190000所以需要先检查每个item的字符数如果超过190000则需要分割
    # 这个功能现在先不做，后续再加
    
    for txt_path in txt_files:
        temp_list = []
        docs = []
        # 如果文件不存在或发生其他错误，将打印错误信息并终止程序。
        try:
            with open(txt_path, 'r', encoding='utf-8') as f:
                txt_contents = f.read()
                # 按照换行符 '\n' 分割文件内容
                temp_list = txt_contents.split('\n\n')
        except FileNotFoundError:
            print(f"错误: 文件 '{txt_path}' 不存在。")
            history.append(["", f"错误: 文件 '{txt_path}' 不存在。"])
            continue
        except Exception as e:
            print(f"发生未知错误: {e}")
            history.append(["", f"发生未知错误: {e}"])
            continue

        file_name = os.path.basename(txt_path)
        # 计算每一个item里面字符数的平均值
        total_chars = 0
        for item in temp_list:
            total_chars += len(item)
        avg_chars = total_chars // len(temp_list)
        print(f"每个item的平均字符数: {avg_chars}")

        # 合并N行内容到一个content中
        merged_docs = []
        for i in range(0, len(temp_list), split_lines):
            merged_content = ""
            for j in range(i, min(i + split_lines, len(temp_list))):
                merged_content += temp_list[j] + "\n"
            if merged_content and merged_content.strip() != "":
                merged_docs.append({
                    "file_path": txt_path,
                    "content": merged_content.strip(),
                })
            else:
                continue
        docs.extend(merged_docs)
        # Saving result to file for later
        path = os.path.join(txt_json_path, file_name+".json")
        with open(path, 'w', encoding='utf-8') as f:
            json.dump(docs, f, ensure_ascii=False, indent=4)
        print(f"{file_name} 文件内容提取完成，已保存到文件: {path}")

def extract_pptx_contents(pptx_files,pptx_json_path):
    # 创建输出文件夹
    if not os.path.exists(pptx_json_path):
        os.makedirs(pptx_json_path)
        print(f"创建输出文件夹: {pptx_json_path}")
    output_pptx_images = os.path.join(pptx_json_path,'extracted_pptx_images')
    if not os.path.exists(output_pptx_images):
        os.makedirs(output_pptx_images)
    for pptx_file in pptx_files:
        ppt_file = []
        ppt = Presentation(pptx_file)
        # 从文件路径中获取文件名
        pptx_file_name = os.path.basename(pptx_file)
        # 遍历每一页幻灯片
        for slide_number, slide in enumerate(ppt.slides):
            ppt_page = {}
            ppt_content=""
            table_text = ""
            notes_text = ""
            ppt_imgs_path=[]
            image_index = 1  # 初始化图片索引
            for shape in slide.shapes:
                print(shape.shape_type)
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    # print(text_frame.text)
                    for paragraph in text_frame.paragraphs:	#text_frame.paragraphs 获取段落
                        print(paragraph.text)
                        ppt_content+=paragraph.text+"\n"
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:  # MSO_SHAPE_TYPE.PICTURE表示图片
                    image = shape.image
                    image_bytes = image.blob
                    image_filename = f'{output_pptx_images}/slide_{slide_number + 1}_image_{image_index}.{image.ext}'
                    with open(image_filename, 'wb') as img_file:
                        img_file.write(image_bytes)
                    print(f"图片 {image_index} 提取完成！")
                    ppt_imgs_path.append(image_filename)
                    image_index += 1  # 更新图片索引
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = [cell.text for cell in row.cells]
                        print(f"表格行: {' | '.join(row_text)}")
                        table_text += ' | '.join(row_text) + '\n'
                    if not table.rows:
                        print("没有表格内容")
            # 提取备注
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                notes_text_frame = notes_slide.notes_text_frame
                for paragraph in notes_text_frame.paragraphs:
                    print(f"备注: {paragraph.text}")
                    notes_text += paragraph.text + '\n'
            else:
                print("没有备注内容")
            ppt_content += table_text +"\n"+ notes_text
            ppt_page['file_path'] = pptx_file
            ppt_page['content'] = ppt_content
            ppt_page['images'] = ppt_imgs_path
            ppt_file.append(ppt_page)
        # 保存提取的内容到json文件
        pptx_json_path = os.path.join(pptx_json_path, pptx_file_name+".json")
        with open(pptx_json_path, 'w', encoding='utf-8') as f:
            json.dump(ppt_file, f, ensure_ascii=False, indent=4)
        print(f"PPT内容提取完成，已保存到文件: {pptx_json_path}")

def pptx_write_to_graphRAG(history,pptx_json_path):
    if not os.path.isdir(os.path.join(os.getcwd(), pptx_json_path)):
        print(f"错误: '{os.path.join(os.getcwd(), pptx_json_path)}' 不是一个有效的目录。")
        history.append(["", f"错误: '{os.path.join(os.getcwd(), pptx_json_path)}' 不是一个有效的目录。"])
        return history

    all_json_contents = {}
    for file_name in os.listdir(os.path.join(os.getcwd(), pptx_json_path)):
        if file_name.lower().endswith('.json'):
            file_path = os.path.join(os.getcwd(), pptx_json_path, file_name)
            # 如果文件不存在或发生其他错误，将打印错误信息并终止程序。
            try:
                # 打开并加载 JSON 文件
                with open(file_path, 'r', encoding='utf-8') as f:
                    json_contents = json.load(f)
                    all_json_contents[file_name] = json_contents
            except FileNotFoundError:
                print(f"错误: 文件 '{file_path}' 不存在。")
                history.append(["", f"错误: 文件 '{file_path}' 不存在。"])
                return history
            except json.JSONDecodeError as e:
                print(f"错误: JSON解析失败 - {e}")
                history.append(["", f"错误: JSON解析失败 - {e}"])
                return history
            except Exception as e:
                print(f"发生未知错误: {e}")
                history.append(["", f"发生未知错误: {e}"])
                return history

    index = PropertyGraphIndex.from_existing(
        property_graph_store=graph_store,
        llm=LlamaOpenAI(model="gpt-4o-mini", temperature=0),
        embed_model=OpenAIEmbedding(model_name="text-embedding-ada-002"),
    )

    total_start_time = time.time()
    # 提取json内容返回整理后的数据list
    for file_name, json_contents in all_json_contents.items():
        # 测量函数执行时间
        start_time = time.time()
        documents_result = pptx_adjust_split_json_contents(json_contents)
        # 调用并行插入函数，默认使用N个并行线程
        parallel_insert(documents_result)
        print("====================Failed sections processing again====================")
        if failed_sections:
            print("Failed to insert the following sections:")
            for document in failed_sections:
                try:
                    # 把内容插入到 Graph 数据库中
                    index.insert(document)
                    print(f"Inserted document: {document.text[:50]}...")
                except Exception as e:
                    print(f"Error inserting document: {e}  section: {document.text[:80]}...")
        # 计算并打印执行时间
        execution_time = time.time() - start_time
        print(f"{file_name}数据导入neo4j执行时间: {round(execution_time, 2)} 秒")
        print(f"{file_name}数据导入neo4j完成")
    # 计算并打印执行时间
    execution_time = time.time() - total_start_time
    print(f"全部数据导入neo4j执行时间: {round(execution_time, 2)} 秒")
    print("全部数据导入neo4j完成")
    history.append(["", f"全部PPT数据导入完成。全部数据导入执行时间: {round(execution_time, 2)} 秒"])
    return history

def extract_xlsx_contents(xlsx_files):
    all_sheet_data = {}
    for xlsx_file in xlsx_files:
        # 获取文件名
        file_name = os.path.basename(xlsx_file)
        # 打开工作簿
        book = openpyxl.load_workbook(xlsx_file)
        # 获取工作表数量和名称
        print("The number of worksheets is {0}".format(len(book.sheetnames)))
        print("Worksheet name(s): {0}".format(book.sheetnames))
        # 创建一个字典来存储所有工作表的数据
        sheet_data = {}
        # 遍历每个工作表
        for sheet_name in book.sheetnames:
            sh = book[sheet_name]
            sheet_content = {}
            
            # 遍历工作表的所有行和列
            for row_idx, row in enumerate(sh.iter_rows(values_only=True), start=1):
                for col_idx, cell_value in enumerate(row, start=1):
                    if cell_value is not None:
                        if isinstance(cell_value, str):
                            cell_value = cell_value.replace('\u3000', '')
                        sheet_content[(row_idx, col_idx)] = cell_value
            # 将工作表的数据存储到字典中
            sheet_data[sheet_name] = sheet_content
        print("========file_name:",file_name)
        print(sheet_data)
        all_sheet_data[file_name] = sheet_data
    return all_sheet_data

def excel_adjust_split_sheet_contents(file_name,sheet_contents):
    documents_result = []
    idx=0
    for sheet_name, sheet_content in sheet_contents.items():
        # 将 sheet_content 中的所有反斜杠 \ 替换为正斜杠 /
        sheet_content = str(sheet_content).strip().replace('\\', '/')
        documents_result.append(Document(text=sheet_content, id_=f"doc_id_{idx}", metadata={"file_name": file_name, "sheet_name": sheet_name}))
        print(f"file_name: {file_name}")
        print(f"sheet_name: {sheet_name}")
        print("_id:",f"doc_id_{idx}")
        idx+=1
    print("==========len documents_result:",len(documents_result))
    return documents_result

def xlsx_write_to_graphRAG(history,all_sheet_data):
    if len(all_sheet_data)==0:
        print("Error: No excel data to process")
        history.append(["", "Error: No excel data to process"])
        return history

    index = PropertyGraphIndex.from_existing(
        property_graph_store=graph_store,
        llm=LlamaOpenAI(model="gpt-4o-mini", temperature=0),
        embed_model=OpenAIEmbedding(model_name="text-embedding-ada-002"),
    )

    total_start_time = time.time()
    # 提取json内容返回整理后的数据list
    for file_name, sheet_contents in all_sheet_data.items():
        # 测量函数执行时间
        start_time = time.time()
        documents_result = excel_adjust_split_sheet_contents(file_name,sheet_contents)
        # 调用并行插入函数，默认使用N个并行线程
        parallel_insert(documents_result)
        print("====================Failed sections processing again====================")
        if failed_sections:
            print("Failed to insert the following sections:")
            for document in failed_sections:
                try:
                    # 把内容插入到 Graph 数据库中
                    index.insert(document)
                    print(f"Inserted document: {document.text[:50]}...")
                except Exception as e:
                    print(f"Error inserting document: {e}  section: {document.text[:80]}...")
        # 计算并打印执行时间
        execution_time = time.time() - start_time
        print(f"{file_name}数据导入neo4j执行时间: {round(execution_time, 2)} 秒")
        print(f"{file_name}数据导入neo4j完成")
    # 计算并打印执行时间
    execution_time = time.time() - total_start_time
    print(f"全部数据导入neo4j执行时间: {round(execution_time, 2)} 秒")
    print("全部数据导入neo4j完成")
    history.append(["", f"全部Excel数据导入完成。全部数据导入执行时间: {round(execution_time, 2)} 秒"])
    return history

def extract_file_paths(data):
    file_paths = []
    img_paths = []
    try:
        # 提取所有 file_path
        file_paths = [info['file_path'] for info in data.values()]
        # 提取所有 image_path
        img_paths = [img for info in data.values() if 'images' in info and info['images'] for img in info['images']]
        # 去除重复的路径（如果需要）
        unique_file_paths = list(set(file_paths))
        unique_img_paths = list(set(img_paths))
    except Exception as e:
        print(f"Error in extracting file paths: {e}")
        unique_file_paths = []
        unique_img_paths = []
    return unique_file_paths,unique_img_paths

# 生成最终回答结果
def generation_answer(history,query_content,nodes_info,file_paths_list):
    prompt_template = """
        你是一个知识回答助理，你根据提供的知识图谱(Graph Insights)检索结果，上下文信息以及问题(query)综合考虑后回复最佳结果，需要使用中文回答问题。

        Graph Insights: {graph_insights}
        Query: {query}

        1. 使用提供的信息(graph_insights)以及上下文信息来回答用户的问题。如果你不知道答案，就说你不知道，不要试图编造答案。
        2. 提供的信息(graph_insights)以及上下文信息不需要全部使用可能有和问题不相关的内容，你需要根据问题的需要选择合适的信息。不相关的信息可能会降低回答的质量。
        3. 知识图谱(Graph Insights)获取的重要方法和信息应该原样保留。不可以省略或改变。
        4. 回复应该使用便于阅读的格式输出。输出格式应该有段落，数字小标，标点符号和适当的换行。
    """
    prompt_template = prompt_template.format(graph_insights=nodes_info, query=query_content)
    # print("=============================",prompt_template)
    history_openai_format = []
    # print("history=====",history)
    for human, assistant in history:
        # 检查每个子列表中的两个元素是否都是字符串
        if isinstance(human, str) and isinstance(assistant, str) and "![image]" not in assistant[:10]:
            history_openai_format.append({"role": "user", "content": human })
            history_openai_format.append({"role": "assistant", "content": assistant})
        else:
            continue
    history_openai_format.append({"role": "user", "content": prompt_template})
    # print("history_openai_format=====",history_openai_format)

    response = client.chat.completions.create(
    model="o1-mini-2024-09-12",
    messages=history_openai_format,
    # temperature=0,
    # max_tokens=16383,
    # top_p=1,
    # frequency_penalty=0,
    # presence_penalty=0,
    # response_format={
    #     "type": "text"
    # }
    )
    # 结果输出
    print(response.choices[0].message.content)
    return response.choices[0].message.content

def only_query(history,query_content):
    # 测量函数执行时间
    start_time = time.time()
    query_engine = index_query.as_query_engine(
        include_text=True,
        similarity_top_k=15,
    )
    # 这个是直接回复内容，可能不够准确应该只返回node节点信息然后交给LLM模型进行回答
    response = query_engine.query(query_content)
    # print("========",str(response))
    file_paths_list, img_paths_list = extract_file_paths(response.metadata)
    # 打印结果
    for result in file_paths_list:
        print(f"文件路径: {result}")
    for result in img_paths_list:
        print(f"image路径: {result}")
    # 这个是返回节点方式 返回一个response列表 然后让LLM模型进行回答
    nodes = query_engine.retrieve(query_content)
    nodes_info=""
    for node in nodes:
        nodes_info+=node.text+"\n"
    # print("\n================================Retrieved nodes:============================\n")
    # print(nodes_info)
    # 去掉None的文件路径
    file_paths_list_filtered = [file for file in file_paths_list if file is not None]
    img_paths_list_filtered = [file for file in img_paths_list if file is not None]
    temp="相关参考文件链接：\n"+"\n".join(file_paths_list_filtered)
    # temp_img="相关参考图片链接：\n"+"\n".join(img_paths_list_filtered)
    history.append([query_content, generation_answer(history,query_content,nodes_info,file_paths_list_filtered)+"\n\n\n"+temp])
    # 转换图片路径为 base64 编码
    base64_images = [img_to_base64(img_path) for img_path in img_paths_list_filtered]
    # 添加图片链接到历史记录
    for base64_image in base64_images:
        history.append(["", f"![image]({base64_image})"])
    # 计算并打印执行时间
    execution_time = time.time() - start_time
    print(f"Query执行时间: {round(execution_time, 2)} 秒")
    print("==================")
    return history

def upload_files(files_path):
    """
    读取指定路径下的多个文件，并返回文件path。
    """
    global upload_files_path
    files_content_temp=""
    print('上传文件的地址：{}'.format(files_path))  # 输出上传后的文件在gradio中保存的绝对地址
    for file_path in files_path:
        if not os.path.exists(file_path):
            error_message= f"文件 {file_path} 不存在"
            raise FileNotFoundError(error_message)
            files_content_temp += f"\n======================{os.path.basename(file_path)}======================\n\n{error_message}\n"
            continue
        files_content_temp += "=====" + os.path.basename(file_path) + "\n"
    # 把用户提交的所有文件路径保存到全局变量中
    upload_files_path=files_path
    for upload_file_path in upload_files_path:
        print("upload_files_path={}".format(upload_file_path))
    return files_content_temp

def bot(history):
    yield history

def print_like_dislike(x: gr.LikeData):
    print(x.index, x.value, x.liked)

# 清理所有历史记录
def clear_history(query_content, chatbot):
    global history
    query_content=""
    chatbot=[]
    history=[]
    print("history",history)
    return query_content, chatbot

def clear_path_json_files(json_path):
    # 检查路径是否存在
    if os.path.exists(json_path):
        # 获取路径下所有子文件夹中的 .json 文件
        json_files = glob.glob(os.path.join(json_path, "**", "*.json"), recursive=True)
        # 删除所有 .json 文件
        for file in json_files:
            os.remove(file)
        print(f"{json_path} 路径下的所有的 JSON 文件已删除，包括子文件夹。")
    else:
        print(f"{json_path} 路径不存在。")

def img_to_base64(image_path):
    """
    将图片转换为 base64 编码的字符串
    """
    with open(image_path, "rb") as img_file:
        encoded_string = base64.b64encode(img_file.read()).decode('utf-8')
    return f"data:image/png;base64,{encoded_string}"

def greet(history,query_content,operation_select,language_options,radio_option):
    print("query_content:",query_content)
    print("operation_select:",operation_select)
    print("language_options:",language_options)
    print("radio_option:",radio_option)

    if radio_option == "精度优先":
        model_type = "o1-mini-2024-09-12"  # 使用精度优先的模型
    elif radio_option == "速度优先":
        model_type = "gpt-4o-mini"  # 使用速度优先的模型

    # 确保 history 是一个列表
    if history is None:
        history = []

    if operation_select == "提交文件及数据处理":
        if not upload_files_path:
            history.append(["", "没有选择要处理的文件。请先选择要处理的文件，可以多个文件同时处理。"])
            print("没有选择要处理的文件。请先选择要处理的文件，可以多个文件同时处理。")
            yield history
        else:
            pdf_files = [file for file in upload_files_path if file.lower().endswith('.pdf')]
            docx_files = [file for file in upload_files_path if file.lower().endswith('.docx')]
            xlsx_files = [file for file in upload_files_path if file.lower().endswith('.xlsx')]
            pptx_files = [file for file in upload_files_path if file.lower().endswith('.pptx')]
            txt_files = [file for file in upload_files_path if file.lower().endswith('.txt')]
            py_files = [file for file in upload_files_path if file.lower().endswith('.py')]
            if len(pdf_files) > 0:
                pdf_json_path = r"D:\Sparticle\llamaindex\parsed_pdf_docs"
                clear_path_json_files(pdf_json_path)
                extract_pdf_contents(pdf_files,pdf_json_path)
                history=pdf_write_to_graphRAG(history,pdf_json_path)
            else:
                print("没有找到要处理的PDF文件。跳过PDF文件处理。")
                history.append(["", "没有找到要处理的PDF文件。跳过PDF文件处理。"])
            if len(docx_files) > 0:
                docx_json_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),parsed_docx_json)
                clear_path_json_files(docx_json_path)
                extract_docx_contents(docx_files,docx_json_path,split_lines=10)
                history=docx_write_to_graphRAG(history,docx_json_path)
            else:
                print("没有找到要处理的Word文件。跳过Word文件处理。")
                history.append(["", "没有找到要处理的Word文件。跳过Word文件处理。"])
            if len(pptx_files) > 0:
                pptx_json_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),parsed_pptx_json)
                clear_path_json_files(pptx_json_path)
                extract_pptx_contents(pptx_files,pptx_json_path)
                history=pptx_write_to_graphRAG(history,pptx_json_path)
            else:
                print("没有找到要处理的PPT文件。跳过PPT文件处理。")
                history.append(["", "没有找到要处理的PPT文件。跳过PPT文件处理。"])
            if len(xlsx_files) > 0:
                all_sheet_data = extract_xlsx_contents(xlsx_files)
                history=xlsx_write_to_graphRAG(history,all_sheet_data)
            else:
                print("没有找到要处理的Excel文件。跳过Excel文件处理。")
                history.append(["", "没有找到要处理的Excel文件。跳过Excel文件处理。"])
            if len(txt_files) > 0:
                txt_json_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),parsed_txt_json)
                clear_path_json_files(txt_json_path)
                extract_txt_contents(txt_files,txt_json_path,split_lines=10)
                history=txt_write_to_graphRAG(history,txt_json_path)
            else:
                print("没有找到要处理的Txt文件。跳过Txt文件处理。")
                history.append(["", "没有找到要处理的Txt文件。跳过Txt文件处理。"])
            if len(py_files) > 0:
                py_json_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),parsed_py_json)
                clear_path_json_files(py_json_path)
                extract_txt_contents(py_files,py_json_path,split_lines=10)
                history=txt_write_to_graphRAG(history,py_json_path)
            else:
                print("没有找到要处理的Py文件。跳过Py文件处理。")
                history.append(["", "没有找到要处理的Py文件。跳过Py文件处理。"])
            yield history
    elif operation_select == "Chatbot提问获取最佳答案" and model_type == "gpt-4o-mini":
        # 测量函数执行时间
        start_time = time.time()
        query_engine = index_query.as_query_engine(
            include_text=True,
            similarity_top_k=15,
        )
        response = query_engine.query(query_content)
        file_paths_list, img_paths_list = extract_file_paths(response.metadata)
        # 打印结果
        for result in file_paths_list:
            print(f"文件路径: {result}")
        # 这个是返回节点方式 返回一个response列表 然后让LLM模型进行回答
        nodes = query_engine.retrieve(query_content)
        nodes_info=""
        for node in nodes:
            nodes_info+=node.text+"\n"
        # print("\n==================================Retrieved nodes:==============================\n")
        # print(nodes_info)
        # 去掉None的文件路径
        file_paths_list_filtered = [file for file in file_paths_list if file is not None]
        img_paths_list_filtered = [file for file in img_paths_list if file is not None]
        temp="相关参考文件链接：\n"+"\n".join(file_paths_list_filtered)
        # temp_img="相关参考图片链接：\n"+"\n".join(img_paths_list_filtered)
        prompt_template = """
        你是一个知识回答助理，你根据提供的知识图谱(Graph Insights)检索结果，上下文信息以及问题(query)综合考虑后回复最佳结果，需要使用中文回答问题。

        Graph Insights: {graph_insights}
        Query: {query}

        1. 使用提供的信息(graph_insights)以及上下文信息来回答用户的问题。如果你不知道答案，就说你不知道，不要试图编造答案。
        2. 提供的信息(graph_insights)以及上下文信息不需要全部使用可能有和问题不相关的内容，你需要根据问题的需要选择合适的信息。不相关的信息可能会降低回答的质量。
        3. 知识图谱(Graph Insights)获取的重要方法和信息应该原样保留。不可以省略或改变。
        4. 回复应该使用便于阅读的格式输出。输出格式应该有段落，数字小标，标点符号和适当的换行。
        """

        prompt_template = prompt_template.format(graph_insights=nodes_info, query=query_content)
        # 转换图片路径为 base64 编码
        base64_images = [img_to_base64(img_path) for img_path in img_paths_list_filtered]
        history_openai_format = []
        for human, assistant in history:
            if isinstance(human, str) and isinstance(assistant, str) and "![image]" not in assistant[:10]:
                history_openai_format.append({"role": "user", "content": human })
                history_openai_format.append({"role": "assistant", "content": assistant})
            else:
                continue
        history_openai_format.append({"role": "user", "content": prompt_template})
        # print("history_openai_format=====",history_openai_format)

        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=history_openai_format,
            temperature=0,
            stream=True
        )
        # response.choices[0].message.content
        partial_message = ""
        for chunk in response:
            if chunk.choices[0].delta.content is not None:
                partial_message += chunk.choices[0].delta.content
                yield history + [[query_content, partial_message]]  # 确保返回的格式为包含两个元素的列表

        history += [[query_content, partial_message+"\n\n"+temp]]  # 确保返回的格式为包含两个元素的列表
        # 添加图片链接到历史记录
        for base64_image in base64_images:
            history.append(["", f"![image]({base64_image})"])
        yield history
        # 计算并打印执行时间
        execution_time = time.time() - start_time
        print(f"Query执行时间: {round(execution_time, 2)} 秒")
        print("====================================")
    elif operation_select == "Chatbot提问获取最佳答案" and model_type == "o1-mini-2024-09-12":
        history=only_query(history,query_content)
        yield history


with tempfile.TemporaryDirectory(dir='.') as tmpdir:
    main_process_input = gr.Interface(fn=upload_files,inputs=gr.File(elem_id="file_upload_component",file_count="multiple",file_types=[".pdf", ".docx", ".xlsx", ".pptx", ".txt", ".py"]), outputs="text")

# gradio内容输出控制
with gr.Blocks() as demo:
    gr.Markdown("""
                <h1 style="text-align: center;">Chatbot AI System</h1>
                """)  # 设置标题 可以使用markdown语法
    chatbot = gr.Chatbot(
        [],
        elem_id="chatbot",
        bubble_full_width=False,
        show_copy_button=True,
        height=650,
        avatar_images=(None, (os.path.join(os.path.dirname(__file__), "icon.png"))),
    )
    language_options = gr.Dropdown(['Chinese', 'Japanese', 'English'], value="Chinese", label="出力言語選択")
    operation_select = gr.Dropdown(['提交文件及数据处理', 'Chatbot提问获取最佳答案'],  value="Chatbot提问获取最佳答案", label="功能选择")
    query_content = gr.Textbox(label="Chatbot提问", lines=1, placeholder="请输入问题", value="扫描工具主要界面包括什么 说明你的理由？")  # 输入框
    radio_option = gr.Radio(["速度优先", "精度优先"], label="模型种类选择", value="速度优先")
    with gr.Row():
        request = gr.Button(value="提交请求")
        btn = gr.Button(value="清空历史记录")
    # 采用tag标签方式实现多个文件输入输出
    gr.TabbedInterface([main_process_input], ["提交文件"])
    btn.click(clear_history, inputs=[query_content, chatbot], outputs=[query_content, chatbot])
    # 通过点击request按钮触发请求
    request.click(greet, inputs=[chatbot,query_content,operation_select,language_options,radio_option],outputs=chatbot)
    chatbot.like(print_like_dislike, None, None)


demo.queue()
PORT="7680"
if __name__ == "__main__":
    demo.launch(share=True)  # 使用auth鉴权参数的时候不要使用share参数
    # demo.launch(server_name="0.0.0.0",server_port=7865,auth=("username", "password"))
    demo.launch(server_name="0.0.0.0",server_port=PORT)
